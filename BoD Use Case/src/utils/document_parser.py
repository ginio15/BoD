"""
Document Parser for BoD Presentation Analysis System

Handles parsing of PDF and PPTX files with OCR fallback for image-based content.
"""

import os
import tempfile
from pathlib import Path
from typing import List, Dict, Any, Optional, Union
import logging

logger = logging.getLogger(__name__)

# Document processing imports
import fitz  # PyMuPDF
from pptx import Presentation

# Optional OCR imports
try:
    import pytesseract
    from PIL import Image
    OCR_AVAILABLE = True
except ImportError as e:
    logger.warning(f"OCR dependencies not available: {e}")
    pytesseract = None
    Image = None
    OCR_AVAILABLE = False

# Local imports
import sys
from pathlib import Path
sys.path.append(str(Path(__file__).parent.parent.parent))

from src.models.document import ProcessedDocument, DocumentPage, DocumentMetadata
from config.settings import Config

class DocumentParser:
    """Main document parser class supporting PDF and PPTX formats"""
    
    def __init__(self):
        self.config = Config()
        self._validate_dependencies()
    
    def _validate_dependencies(self):
        """Validate that required dependencies are available"""
        if OCR_AVAILABLE:
            try:
                # Check Tesseract installation
                pytesseract.get_tesseract_version()
                logger.info("Tesseract OCR is available")
            except Exception as e:
                logger.warning(f"Tesseract OCR not available: {e}")
                logger.warning("OCR functionality will be limited")
        else:
            logger.warning("OCR dependencies not installed - OCR functionality disabled")
    
    def process_document(self, file_path: Union[str, Path], use_ocr: bool = True) -> ProcessedDocument:
        """
        Process a document and extract structured content with metadata
        
        Args:
            file_path: Path to the document file
            use_ocr: Whether to apply OCR to image-based content
            
        Returns:
            ProcessedDocument with extracted content and metadata
        """
        file_path = Path(file_path)
        
        if not file_path.exists():
            raise FileNotFoundError(f"Document not found: {file_path}")
        
        # Check file size
        file_size_mb = file_path.stat().st_size / (1024 * 1024)
        if file_size_mb > self.config.DOCUMENT_CONFIG["max_file_size_mb"]:
            raise ValueError(f"File too large: {file_size_mb:.1f}MB > {self.config.DOCUMENT_CONFIG['max_file_size_mb']}MB")
        
        # Determine file type and process
        file_extension = file_path.suffix.lower()
        
        if file_extension == ".pdf":
            return self._process_pdf(file_path, use_ocr)
        elif file_extension == ".pptx":
            return self._process_pptx(file_path, use_ocr)
        elif file_extension == ".txt":
            return self._process_text_file(file_path)
        else:
            raise ValueError(f"Unsupported file format: {file_extension}")
    
    def _process_pdf(self, pdf_path: Path, use_ocr: bool) -> ProcessedDocument:
        """Process PDF document"""
        logger.info(f"Processing PDF: {pdf_path.name}")
        
        document = fitz.open(pdf_path)
        pages = []
        total_text = ""
        
        for page_num in range(len(document)):
            page = document[page_num]
            
            # Extract text
            text = page.get_text()
            
            # Calculate text density for OCR decision
            text_density = len(text.strip()) / max(len(text), 1)
            
            # Apply OCR if text density is low and OCR is enabled
            if use_ocr and text_density < self.config.DOCUMENT_CONFIG["text_density_threshold"]:
                try:
                    ocr_text = self._apply_ocr_to_page(page)
                    if len(ocr_text.strip()) > len(text.strip()):
                        text = ocr_text
                        source = "ocr"
                    else:
                        source = "mixed"
                except Exception as e:
                    logger.warning(f"OCR failed for page {page_num + 1}: {e}")
                    source = "text"
            else:
                source = "text"
            
            # Extract page metadata
            page_metadata = self._extract_page_metadata(page, text)
            
            # Create page object
            doc_page = DocumentPage(
                page_number=page_num + 1,
                text=text,
                source=source,
                metadata=page_metadata
            )
            
            pages.append(doc_page)
            total_text += text + "\n"
        
        document.close()
        
        # Create document metadata
        metadata = DocumentMetadata(
            filename=pdf_path.name,
            file_type="pdf",
            total_pages=len(pages),
            word_count=len(total_text.split()),
            char_count=len(total_text),
            file_size_mb=pdf_path.stat().st_size / (1024 * 1024)
        )
        
        return ProcessedDocument(
            pages=pages,
            metadata=metadata,
            full_text=total_text
        )
    
    def _process_pptx(self, pptx_path: Path, use_ocr: bool) -> ProcessedDocument:
        """Process PowerPoint presentation"""
        logger.info(f"Processing PPTX: {pptx_path.name}")
        
        presentation = Presentation(pptx_path)
        pages = []
        total_text = ""
        
        for slide_num, slide in enumerate(presentation.slides):
            # Extract text from slide
            text = self._extract_slide_text(slide)
            
            # Extract tables if present
            tables = self._extract_slide_tables(slide)
            if tables:
                text += "\n" + self._format_tables_as_text(tables)
            
            # Check if OCR is needed for images
            source = "text"
            if use_ocr and self._slide_has_images(slide):
                try:
                    # Extract slide as image and apply OCR
                    ocr_text = self._apply_ocr_to_slide(slide, slide_num)
                    if ocr_text and len(ocr_text.strip()) > len(text.strip()) * 1.2:
                        text += "\n" + ocr_text
                        source = "mixed"
                except Exception as e:
                    logger.warning(f"OCR failed for slide {slide_num + 1}: {e}")
            
            # Extract slide metadata
            slide_metadata = {
                "slide_title": self._extract_slide_title(slide),
                "shape_count": len(slide.shapes),
                "has_tables": bool(tables),
                "has_images": self._slide_has_images(slide)
            }
            
            # Create page object
            doc_page = DocumentPage(
                page_number=slide_num + 1,
                text=text,
                source=source,
                metadata=slide_metadata
            )
            
            pages.append(doc_page)
            total_text += text + "\n"
        
        # Create document metadata
        metadata = DocumentMetadata(
            filename=pptx_path.name,
            file_type="pptx",
            total_pages=len(pages),
            word_count=len(total_text.split()),
            char_count=len(total_text),
            file_size_mb=pptx_path.stat().st_size / (1024 * 1024)
        )
        
        return ProcessedDocument(
            pages=pages,
            metadata=metadata,
            full_text=total_text
        )
    
    def _process_text_file(self, text_path: Path) -> ProcessedDocument:
        """Process plain text file"""
        logger.info(f"Processing text file: {text_path.name}")
        
        # Read file content
        with open(text_path, 'r', encoding='utf-8') as f:
            content = f.read()
        
        # Create single page from text content
        page = DocumentPage(
            page_number=1,
            text=content,
            source="text",
            metadata={
                "file_type": "text",
                "encoding": "utf-8"
            }
        )
        
        # Create document metadata
        metadata = DocumentMetadata(
            filename=text_path.name,
            file_type="txt",
            total_pages=1,
            word_count=len(content.split()),
            char_count=len(content),
            file_size_mb=text_path.stat().st_size / (1024 * 1024)
        )
        
        return ProcessedDocument(
            pages=[page],
            metadata=metadata,
            full_text=content
        )
    
    def _apply_ocr_to_page(self, page) -> str:
        """Apply OCR to a PDF page"""
        if not OCR_AVAILABLE:
            logger.warning("OCR not available - skipping OCR for page")
            return ""
        
        try:
            # Convert page to image
            pix = page.get_pixmap()
            img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
            
            # Apply OCR
            ocr_config = self.config.OCR_CONFIG["config"]
            text = pytesseract.image_to_string(img, config=ocr_config)
            
            return text.strip()
        except Exception as e:
            logger.error(f"OCR failed for page: {e}")
            return ""
    
    def _apply_ocr_to_slide(self, slide, slide_num: int) -> str:
        """Apply OCR to a PowerPoint slide by converting to image"""
        if not OCR_AVAILABLE:
            logger.warning("OCR not available - skipping slide OCR")
            return ""
            
        try:
            # For PPTX slides with images (like screenshots), we need to:
            # 1. Extract images from the slide
            # 2. Apply OCR to those images
            
            ocr_text_parts = []
            
            # Look for image shapes in the slide
            for shape in slide.shapes:
                if shape.shape_type == 13:  # Picture shape type
                    try:
                        # Extract image data from shape
                        image_ocr = self._extract_and_ocr_image_from_shape(shape, slide_num)
                        if image_ocr:
                            ocr_text_parts.append(image_ocr)
                    except Exception as e:
                        logger.warning(f"Failed to OCR image in slide {slide_num + 1}: {e}")
                        continue
            
            # If no images were processed, try alternative approach
            if not ocr_text_parts:
                logger.debug(f"No OCR-able images found in slide {slide_num + 1}")
                return ""
            
            combined_text = "\n".join(ocr_text_parts)
            logger.info(f"Extracted {len(combined_text)} characters via OCR from slide {slide_num + 1}")
            return combined_text
            
        except Exception as e:
            logger.error(f"OCR failed for slide {slide_num + 1}: {e}")
            return ""
    
    def _extract_and_ocr_image_from_shape(self, shape, slide_num: int) -> str:
        """Extract image from shape and apply OCR optimized for presentations"""
        try:
            # Get the image data from the shape
            image_part = shape.image.blob
            
            # Create PIL Image from bytes
            import io
            image = Image.open(io.BytesIO(image_part))
            
            # Preprocess image for better OCR results
            processed_image = self._preprocess_image_for_ocr(image)
            
            # Apply OCR with presentation-optimized settings
            ocr_config = self.config.OCR_CONFIG.get("presentation_config", self.config.OCR_CONFIG.get("config", "--psm 6"))
            text = pytesseract.image_to_string(processed_image, config=ocr_config)
            
            # Clean and validate the extracted text
            cleaned_text = self._clean_ocr_text(text)
            
            return cleaned_text
            
        except Exception as e:
            logger.warning(f"Failed to extract image from shape in slide {slide_num + 1}: {e}")
            return ""
    
    def _preprocess_image_for_ocr(self, image):
        """Preprocess image to improve OCR accuracy for presentation screenshots"""
        try:
            # Convert to RGB if necessary
            if image.mode != 'RGB':
                image = image.convert('RGB')
            
            # Convert to grayscale for better OCR
            image = image.convert('L')
            
            # Enhanced contrast for presentation screenshots
            from PIL import ImageEnhance, ImageFilter
            contrast_factor = self.config.OCR_CONFIG.get("preprocessing", {}).get("enhance_contrast", 1.8)
            enhancer = ImageEnhance.Contrast(image)
            image = enhancer.enhance(contrast_factor)
            
            # Sharpening for better text clarity
            image = image.filter(ImageFilter.SHARPEN)
            
            # Resize for optimal OCR (presentations often need larger sizes)
            width, height = image.size
            min_width = self.config.OCR_CONFIG.get("preprocessing", {}).get("min_resize_width", 1200)
            min_height = self.config.OCR_CONFIG.get("preprocessing", {}).get("min_resize_height", 800)
            
            if width < min_width or height < min_height:
                scale_factor = max(min_width / width, min_height / height)
                new_width = int(width * scale_factor)
                new_height = int(height * scale_factor)
                image = image.resize((new_width, new_height), Image.LANCZOS)
                logger.debug(f"Resized image from {width}x{height} to {new_width}x{new_height}")
            
            return image
            
        except Exception as e:
            logger.warning(f"Image preprocessing failed: {e}")
            return image  # Return original if preprocessing fails
    
    def _clean_ocr_text(self, text: str) -> str:
        """Clean and validate OCR extracted text from presentations"""
        if not text:
            return ""
        
        # Remove excessive whitespace and normalize line breaks
        cleaned = " ".join(text.split())
        
        # Common OCR corrections for presentation content
        corrections = {
            'O%': '0%',  # Common OCR mistake
            'I%': '1%',
            'S%': '5%',
            ' %': '%',  # Fix spacing before percentages
            '% ': '% ',  # Normalize percentage spacing
            '( ': '(',  # Fix parentheses spacing
            ' )': ')',
            'H/A': 'H/A',  # Preserve specific terms
        }
        
        for wrong, right in corrections.items():
            cleaned = cleaned.replace(wrong, right)
        
        # Enhanced word separation for concatenated OCR text
        cleaned = self._separate_concatenated_words(cleaned)
        
        # Remove very short "words" that are likely OCR noise (but keep numbers and percentages)
        words = cleaned.split()
        filtered_words = []
        for word in words:
            # Keep if: length >= 2, or is digit, or contains %, or is common short word
            if (len(word) >= 2 or 
                word.isdigit() or 
                '%' in word or 
                word.lower() in {'a', 'i', 'q1', 'q2', 'q3', 'q4', 'h/a', 'gr'}):
                filtered_words.append(word)
        
        # Only return text if we have substantial content
        result = " ".join(filtered_words)
        if len(result) < 8:  # Lower threshold for presentations (they can be concise)
            logger.debug("OCR text too short, likely noise")
            return ""
        
        return result

    def _separate_concatenated_words(self, text: str) -> str:
        """Separate concatenated words in OCR text using common patterns"""
        import re
        
        # Split on capital letters followed by lowercase (CamelCase)
        # But preserve abbreviations and percentages
        result = text
        
        # Split CamelCase: "ProjectKalamaras" -> "Project Kalamaras"
        result = re.sub(r'([a-z])([A-Z])', r'\1 \2', result)
        
        # Split before numbers: "Phase1" -> "Phase 1", "YE2024" -> "YE 2024"  
        result = re.sub(r'([a-zA-Z])(\d)', r'\1 \2', result)
        
        # Split after numbers: "2024Target" -> "2024 Target", "5mil" -> "5 mil"
        result = re.sub(r'(\d)([a-zA-Z])', r'\1 \2', result)
        
        # Split on common business terms boundaries
        business_terms = [
            'Target', 'Status', 'Initiatives', 'Overview', 'Update', 'Strategy', 
            'Negotiations', 'Sourcing', 'Inventory', 'Reduction', 'Resizing',
            'Staffing', 'Linking', 'Automation', 'System', 'Phase', 'Delivered',
            'Ongoing', 'Track', 'Started', 'Implementation', 'Work', 'Working',
            'Week', 'Associates', 'Running', 'Enable', 'Delivery', 'Planned',
            'Order', 'Network', 'Replenishment', 'Department', 'Management'
        ]
        
        for term in business_terms:
            # Split when term appears concatenated
            pattern = f'([a-z]){term}'
            result = re.sub(pattern, f'\\1 {term}', result, flags=re.IGNORECASE)
            
            pattern = f'{term}([A-Z][a-z])'
            result = re.sub(pattern, f'{term} \\1', result, flags=re.IGNORECASE)
        
        # Split on percentage boundaries: "Ent:55%Sourcing" -> "Ent:55% Sourcing"
        result = re.sub(r'(%\s*)([A-Z][a-z])', r'\1 \2', result)
        
        # Split on colon boundaries: "Linking:Ongoing" -> "Linking: Ongoing"
        result = re.sub(r'(:)([A-Z][a-z])', r'\1 \2', result)
        
        # Clean up excessive spaces
        result = re.sub(r'\s+', ' ', result).strip()
        
        return result
    
    def _extract_slide_text(self, slide) -> str:
        """Extract text from all text shapes in a slide"""
        text_content = []
        
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text_content.append(shape.text.strip())
        
        return "\n".join(text_content)
    
    def _extract_slide_tables(self, slide) -> List[List[List[str]]]:
        """Extract tables from a slide"""
        tables = []
        
        for shape in slide.shapes:
            if shape.has_table:
                table_data = []
                for row in shape.table.rows:
                    row_data = [cell.text.strip() for cell in row.cells]
                    table_data.append(row_data)
                tables.append(table_data)
        
        return tables
    
    def _format_tables_as_text(self, tables: List[List[List[str]]]) -> str:
        """Format extracted tables as readable text"""
        formatted_tables = []
        
        for i, table in enumerate(tables):
            table_text = f"\n[Table {i+1}]\n"
            for row in table:
                table_text += " | ".join(row) + "\n"
            formatted_tables.append(table_text)
        
        return "\n".join(formatted_tables)
    
    def _slide_has_images(self, slide) -> bool:
        """Check if slide contains images"""
        for shape in slide.shapes:
            if shape.shape_type == 13:  # Picture shape type
                return True
        return False
    
    def _extract_slide_title(self, slide) -> str:
        """Extract title from slide"""
        if slide.shapes.title and slide.shapes.title.text:
            return slide.shapes.title.text.strip()
        
        # Fallback: look for the first text shape
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                return shape.text.strip()[:100]  # First 100 chars as title
        
        return "Untitled Slide"
    
    def _extract_page_metadata(self, page, text: str) -> Dict[str, Any]:
        """Extract metadata from a PDF page"""
        return {
            "page_title": self._extract_page_title(text),
            "text_length": len(text),
            "word_count": len(text.split()),
            "has_images": len(page.get_images()) > 0,
            "has_text": len(text.strip()) > 0
        }
    
    def _extract_page_title(self, text: str) -> str:
        """Extract title from page text"""
        lines = text.strip().split('\n')
        if lines:
            # Return first non-empty line as title
            for line in lines:
                if line.strip():
                    return line.strip()[:100]
        return "Untitled Page"

class DocumentValidator:
    """Validates documents before processing"""
    
    @staticmethod
    def validate_file(file_path: Union[str, Path]) -> Dict[str, Any]:
        """Validate a file for processing"""
        file_path = Path(file_path)
        
        validation_result = {
            "valid": True,
            "errors": [],
            "warnings": [],
            "metadata": {}
        }
        
        # Check if file exists
        if not file_path.exists():
            validation_result["valid"] = False
            validation_result["errors"].append(f"File not found: {file_path}")
            return validation_result
        
        # Check file extension
        if file_path.suffix.lower() not in Config.DOCUMENT_CONFIG["supported_formats"]:
            validation_result["valid"] = False
            validation_result["errors"].append(f"Unsupported format: {file_path.suffix}")
        
        # Check file size
        file_size_mb = file_path.stat().st_size / (1024 * 1024)
        if file_size_mb > Config.DOCUMENT_CONFIG["max_file_size_mb"]:
            validation_result["valid"] = False
            validation_result["errors"].append(f"File too large: {file_size_mb:.1f}MB")
        elif file_size_mb > Config.DOCUMENT_CONFIG["max_file_size_mb"] * 0.8:
            validation_result["warnings"].append(f"Large file: {file_size_mb:.1f}MB")
        
        # Add metadata
        validation_result["metadata"] = {
            "filename": file_path.name,
            "size_mb": file_size_mb,
            "extension": file_path.suffix.lower()
        }
        
        return validation_result
