#!/usr/bin/env python3
"""
Detailed OCR Text Inspection
Examines exactly what text is being extracted and provides detailed analysis
"""

import os
import sys
from pathlib import Path

# Add the src directory to the Python path
sys.path.insert(0, os.path.join(os.path.dirname(__file__), 'src'))

from src.utils.document_parser import DocumentParser
from pptx import Presentation

def inspect_pptx_structure(file_path: str):
    """Inspect the actual structure of a PPTX file"""
    print(f"🔍 PPTX Structure Analysis: {Path(file_path).name}")
    print("=" * 60)
    
    prs = Presentation(file_path)
    
    total_text_shapes = 0
    total_image_shapes = 0
    total_table_shapes = 0
    total_other_shapes = 0
    
    all_text_content = []
    
    for slide_idx, slide in enumerate(prs.slides):
        print(f"\n📊 Slide {slide_idx + 1}:")
        
        slide_text_content = []
        slide_images = 0
        slide_tables = 0
        slide_other = 0
        
        for shape_idx, shape in enumerate(slide.shapes):
            print(f"   Shape {shape_idx + 1}: ", end="")
            
            if hasattr(shape, "text") and shape.text.strip():
                text = shape.text.strip()
                slide_text_content.append(text)
                total_text_shapes += 1
                word_count = len(text.split())
                print(f"TEXT ({word_count} words) - '{text[:50]}{'...' if len(text) > 50 else ''}'")
                
            elif hasattr(shape, "image"):
                slide_images += 1
                total_image_shapes += 1
                print(f"IMAGE - {shape.image.blob[:20] if hasattr(shape.image, 'blob') else 'No blob'}")
                
            elif shape.has_table:
                slide_tables += 1
                total_table_shapes += 1
                print(f"TABLE ({len(shape.table.rows)} rows, {len(shape.table.columns)} cols)")
                
                # Extract table text
                table_text = []
                for row in shape.table.rows:
                    for cell in row.cells:
                        if cell.text.strip():
                            table_text.append(cell.text.strip())
                
                if table_text:
                    table_content = " ".join(table_text)
                    slide_text_content.append(table_content)
                    table_word_count = len(table_content.split())
                    print(f" - Table text: {table_word_count} words")
                    
            else:
                slide_other += 1
                total_other_shapes += 1
                print(f"OTHER (type: {type(shape).__name__})")
        
        slide_all_text = " ".join(slide_text_content)
        slide_word_count = len(slide_all_text.split()) if slide_all_text else 0
        
        print(f"   📝 Slide {slide_idx + 1} Summary:")
        print(f"      Text shapes: {len([t for t in slide_text_content if t])} | Images: {slide_images} | Tables: {slide_tables} | Other: {slide_other}")
        print(f"      Total words in slide: {slide_word_count}")
        
        all_text_content.extend(slide_text_content)
    
    # Summary
    total_manual_text = " ".join(all_text_content)
    manual_word_count = len(total_manual_text.split()) if total_manual_text else 0
    
    print(f"\n📊 MANUAL EXTRACTION SUMMARY:")
    print(f"   Total shapes: {total_text_shapes + total_image_shapes + total_table_shapes + total_other_shapes}")
    print(f"   Text shapes: {total_text_shapes}")
    print(f"   Image shapes: {total_image_shapes}")
    print(f"   Table shapes: {total_table_shapes}")
    print(f"   Other shapes: {total_other_shapes}")
    print(f"   Manual word count: {manual_word_count}")
    print(f"   Manual char count: {len(total_manual_text)}")
    
    return {
        "manual_text": total_manual_text,
        "manual_word_count": manual_word_count,
        "manual_char_count": len(total_manual_text),
        "text_shapes": total_text_shapes,
        "image_shapes": total_image_shapes,
        "table_shapes": total_table_shapes
    }

def inspect_ocr_extraction(file_path: str):
    """Inspect what our OCR system extracts"""
    print(f"\n🤖 OCR EXTRACTION ANALYSIS: {Path(file_path).name}")
    print("=" * 60)
    
    parser = DocumentParser()
    processed_doc = parser.process_document(file_path)
    
    print(f"   OCR word count: {len(processed_doc.full_text.split())}")
    print(f"   OCR char count: {len(processed_doc.full_text)}")
    print(f"   OCR pages processed: {len(processed_doc.pages)}")
    
    print(f"\n📄 OCR Extracted Text:")
    print(f"'{processed_doc.full_text}'")
    
    return {
        "ocr_text": processed_doc.full_text,
        "ocr_word_count": len(processed_doc.full_text.split()),
        "ocr_char_count": len(processed_doc.full_text),
        "pages": len(processed_doc.pages)
    }

def compare_results(manual_data: dict, ocr_data: dict, file_name: str):
    """Compare manual extraction vs OCR results"""
    print(f"\n🔄 COMPARISON ANALYSIS: {file_name}")
    print("=" * 60)
    
    word_diff = manual_data["manual_word_count"] - ocr_data["ocr_word_count"]
    char_diff = manual_data["manual_char_count"] - ocr_data["ocr_char_count"]
    
    print(f"📊 Word Count Comparison:")
    print(f"   Manual extraction: {manual_data['manual_word_count']} words")
    print(f"   OCR extraction: {ocr_data['ocr_word_count']} words")
    print(f"   Difference: {word_diff} words ({'Manual higher' if word_diff > 0 else 'OCR higher' if word_diff < 0 else 'Same'})")
    
    print(f"\n📊 Character Count Comparison:")
    print(f"   Manual extraction: {manual_data['manual_char_count']} chars")
    print(f"   OCR extraction: {ocr_data['ocr_char_count']} chars")
    print(f"   Difference: {char_diff} chars")
    
    accuracy = (ocr_data['ocr_word_count'] / manual_data['manual_word_count'] * 100) if manual_data['manual_word_count'] > 0 else 0
    print(f"\n📈 OCR Accuracy: {accuracy:.1f}% of manual count")
    
    # Show text content differences
    manual_words = set(manual_data["manual_text"].lower().split())
    ocr_words = set(ocr_data["ocr_text"].lower().split())
    
    missing_in_ocr = manual_words - ocr_words
    extra_in_ocr = ocr_words - manual_words
    
    if missing_in_ocr:
        print(f"\n⚠️  Words in manual but missing in OCR ({len(missing_in_ocr)}):")
        print(f"   {', '.join(list(missing_in_ocr)[:10])}{'...' if len(missing_in_ocr) > 10 else ''}")
    
    if extra_in_ocr:
        print(f"\n➕ Words in OCR but not in manual ({len(extra_in_ocr)}):")
        print(f"   {', '.join(list(extra_in_ocr)[:10])}{'...' if len(extra_in_ocr) > 10 else ''}")

def main():
    """Main inspection function"""
    print("🚀 Detailed PPTX Text Analysis")
    print("=" * 80)
    
    test_files = [
        "/Users/ginio/projects/Olympia/BoD Use Case/data/uploads/Presentation1.pptx",
        "/Users/ginio/projects/Olympia/BoD Use Case/data/uploads/Presentation2.pptx"
    ]
    
    for file_path in test_files:
        if Path(file_path).exists():
            # Manual structure analysis
            manual_data = inspect_pptx_structure(file_path)
            
            # OCR extraction analysis
            ocr_data = inspect_ocr_extraction(file_path)
            
            # Comparison
            compare_results(manual_data, ocr_data, Path(file_path).name)
            
            print("\n" + "="*80 + "\n")
        else:
            print(f"❌ File not found: {file_path}")

if __name__ == "__main__":
    main()
